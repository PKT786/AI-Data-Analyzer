"""
===========================================================
AI DATA ANALYZER PRO
Dashboard Exporter
Version : 2.0 Enterprise
-----------------------------------------------------------

Responsible for exporting a Dashboard to downloadable files.

Public Methods (API_CONTRACT.md)

    export_pdf(dashboard) -> bytes
    export_png(dashboard) -> bytes
    export_excel(dashboard) -> bytes
    export_powerpoint(dashboard) -> bytes
    export_html(dashboard) -> str

Input: Dashboard. Output: file bytes / str ready for
st.download_button.
===========================================================
"""

from __future__ import annotations

import io
import logging
from datetime import datetime
from typing import Optional

import pandas as pd

from components.dashboard_models import Dashboard

logger = logging.getLogger(__name__)


class DashboardExporter:
    """Exports a Dashboard object to various premium file formats."""

    # ------------------------------------------------------
    # Composed "Dashboard View" image - mirrors the actual
    # app layout: gradient header, KPI cards, chart grid.
    # Used by export_png() and embedded into the Excel /
    # PowerPoint exports so they match what is on screen.
    # ------------------------------------------------------

    def _render_dashboard_image(self, dashboard: Dashboard, columns: int = 3):
        from PIL import Image, ImageDraw, ImageFont

        theme = dashboard.theme
        card_w, chart_h = 620, 380
        title_h = 40
        card_h = title_h + chart_h + 24
        gap = 24
        margin = 30
        header_h = 130
        kpi_h = 110 if dashboard.kpis else 0

        widgets = [w for w in dashboard.widgets if w.visible and w.chart is not None]
        rows = max(1, -(-len(widgets) // columns)) if widgets else 1

        canvas_w = margin * 2 + columns * card_w + (columns - 1) * gap
        canvas_h = (
            header_h + kpi_h + margin
            + rows * card_h + (rows - 1) * gap
            + margin
        )

        canvas = Image.new("RGB", (canvas_w, canvas_h), self._hex_to_rgb255(theme.background))
        draw = ImageDraw.Draw(canvas)

        font_regular, font_bold, font_title, font_small = self._load_fonts()

        # ---- Header banner (gradient) ----
        header_top = self._hex_to_rgb255(theme.primary)
        header_bottom = self._hex_to_rgb255(theme.secondary)
        for y in range(header_h):
            t = y / max(header_h - 1, 1)
            blended = tuple(int(header_top[i] + (header_bottom[i] - header_top[i]) * t) for i in range(3))
            draw.line([(0, y), (canvas_w, y)], fill=blended)

        header_text_color = self._hex_to_rgb255(theme.header_text_color)
        draw.text((margin, 26), dashboard.title, font=font_title, fill=header_text_color)
        draw.text((margin, 74), dashboard.subtitle or "", font=font_regular, fill=header_text_color)

        # ---- KPI cards ----
        y_cursor = header_h + margin
        if dashboard.kpis:
            kpi_count = len(dashboard.kpis)
            kpi_w = (canvas_w - margin * 2 - gap * (kpi_count - 1)) / kpi_count
            for i, kpi in enumerate(dashboard.kpis):
                x0 = margin + i * (kpi_w + gap)
                x1 = x0 + kpi_w
                y1 = y_cursor + kpi_h - 20
                self._rounded_card(draw, (x0, y_cursor, x1, y1), theme)
                draw.text((x0 + 16, y_cursor + 14), str(kpi.get("label", "")).upper(),
                          font=font_small, fill=self._hex_to_rgb255(theme.muted_text))
                draw.text((x0 + 16, y_cursor + 38), str(kpi.get("value", "")),
                          font=font_bold, fill=self._hex_to_rgb255(theme.primary))
            y_cursor += kpi_h

        # ---- Chart grid ----
        for idx, widget in enumerate(widgets):
            row, col = divmod(idx, columns)
            x0 = margin + col * (card_w + gap)
            y0 = y_cursor + row * (card_h + gap)
            x1 = x0 + card_w
            y1 = y0 + card_h

            self._rounded_card(draw, (x0, y0, x1, y1), theme)
            draw.text((x0 + 16, y0 + 10), widget.title, font=font_bold,
                       fill=self._hex_to_rgb255(theme.text_color))

            try:
                png_bytes = widget.chart.to_image(
                    format="png", width=card_w - 32, height=chart_h, scale=2
                )
                chart_img = Image.open(io.BytesIO(png_bytes)).convert("RGB")
                chart_img = chart_img.resize((card_w - 32, chart_h))
                canvas.paste(chart_img, (x0 + 16, y0 + title_h))
            except Exception as ex:
                logger.warning("Could not render widget '%s' for dashboard image: %s", widget.title, ex)

        return canvas

    def export_png(self, dashboard: Dashboard) -> Optional[bytes]:
        """Export a single PNG that mirrors the current on-screen dashboard view."""
        try:
            canvas = self._render_dashboard_image(dashboard)
            buffer = io.BytesIO()
            canvas.save(buffer, format="PNG")
            return buffer.getvalue()
        except Exception as ex:
            logger.exception("PNG export failed: %s", ex)
            return None

    # ------------------------------------------------------
    # PDF - simple, robust: title page + one chart image per page
    # ------------------------------------------------------

    def export_pdf(self, dashboard: Dashboard) -> Optional[bytes]:
        try:
            from reportlab.lib.pagesizes import landscape, A4
            from reportlab.lib.units import cm
            from reportlab.lib.utils import ImageReader
            from reportlab.pdfgen import canvas as pdf_canvas
            from PIL import Image

            buffer = io.BytesIO()
            page_size = landscape(A4)
            c = pdf_canvas.Canvas(buffer, pagesize=page_size)
            width, height = page_size

            # Title page
            c.setFillColorRGB(*self._hex_to_rgb(dashboard.theme.primary))
            c.rect(0, 0, width, height, fill=True, stroke=False)
            c.setFillColorRGB(1, 1, 1)
            c.setFont("Helvetica-Bold", 28)
            c.drawCentredString(width / 2, height / 2 + 20, dashboard.title)
            c.setFont("Helvetica", 14)
            c.drawCentredString(width / 2, height / 2 - 15, dashboard.subtitle or "")
            c.drawCentredString(
                width / 2, height / 2 - 45,
                f"Generated {datetime.now().strftime('%d %b %Y %H:%M')}",
            )
            c.showPage()

            for widget in dashboard.widgets:
                if widget.chart is None:
                    continue
                try:
                    png_bytes = widget.chart.to_image(format="png", width=1200, height=700, scale=2)
                    img = Image.open(io.BytesIO(png_bytes))
                    img_ratio = img.width / img.height
                    draw_width = width - 2 * cm
                    draw_height = draw_width / img_ratio
                    if draw_height > height - 3 * cm:
                        draw_height = height - 3 * cm
                        draw_width = draw_height * img_ratio

                    c.setFillColorRGB(0, 0, 0)
                    c.setFont("Helvetica-Bold", 16)
                    c.drawString(1 * cm, height - 1.2 * cm, widget.title)

                    img_reader_buf = ImageReader(io.BytesIO(png_bytes))
                    c.drawImage(
                        img_reader_buf,
                        (width - draw_width) / 2,
                        (height - draw_height) / 2 - 0.3 * cm,
                        width=draw_width,
                        height=draw_height,
                        preserveAspectRatio=True,
                    )
                    c.showPage()
                except Exception as ex:
                    logger.warning("Skipping widget '%s' in PDF export: %s", widget.title, ex)

            c.save()
            return buffer.getvalue()

        except Exception as ex:
            logger.exception("PDF export failed: %s", ex)
            return None

    # ------------------------------------------------------
    # Excel - KPI sheet + one data sheet per widget
    # ------------------------------------------------------

    def export_excel(self, dashboard: Dashboard) -> Optional[bytes]:
        try:
            buffer = io.BytesIO()
            with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
                workbook = writer.book

                # ---- Dashboard View sheet (matches the on-screen dashboard) ----
                try:
                    image = self._render_dashboard_image(dashboard)
                    img_buffer = io.BytesIO()
                    image.save(img_buffer, format="PNG")
                    img_buffer.seek(0)

                    dash_sheet = workbook.add_worksheet("Dashboard View")
                    writer.sheets["Dashboard View"] = dash_sheet
                    scale = min(1.0, 1400 / image.width)
                    dash_sheet.insert_image(
                        "A1", "dashboard.png",
                        {"image_data": img_buffer, "x_scale": scale, "y_scale": scale},
                    )
                except Exception as ex:
                    logger.warning("Could not embed dashboard image in Excel: %s", ex)

                kpi_rows = [
                    {"KPI": k.get("label"), "Value": k.get("value")}
                    for k in dashboard.kpis
                ]
                pd.DataFrame(kpi_rows or [{"KPI": "N/A", "Value": "N/A"}]).to_excel(
                    writer, sheet_name="KPIs", index=False
                )

                for i, widget in enumerate(dashboard.widgets, start=1):
                    if widget.data is None:
                        continue
                    try:
                        data = widget.data
                        if not isinstance(data, pd.DataFrame):
                            data = pd.DataFrame(data)
                        sheet_name = f"{i}_{widget.title}"[:31].replace("/", "-")
                        data.to_excel(writer, sheet_name=sheet_name, index=False)
                    except Exception as ex:
                        logger.warning("Skipping widget data '%s' in Excel export: %s", widget.title, ex)

            return buffer.getvalue()

        except Exception as ex:
            logger.exception("Excel export failed: %s", ex)
            return None

    # ------------------------------------------------------
    # PowerPoint - one slide per widget, images embedded
    # ------------------------------------------------------

    def export_powerpoint(self, dashboard: Dashboard) -> Optional[bytes]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            blank = prs.slide_layouts[6]

            # Title slide
            slide = prs.slides.add_slide(blank)
            rgb = self._hex_to_rgb(dashboard.theme.primary)
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(int(rgb[0] * 255), int(rgb[1] * 255), int(rgb[2] * 255))

            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5))
            tf = title_box.text_frame
            tf.text = dashboard.title
            tf.paragraphs[0].font.size = Pt(40)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
            tf2 = subtitle_box.text_frame
            tf2.text = dashboard.subtitle or ""
            tf2.paragraphs[0].font.size = Pt(18)
            tf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            # ---- Dashboard Overview slide (matches the on-screen dashboard) ----
            try:
                image = self._render_dashboard_image(dashboard)
                img_buffer = io.BytesIO()
                image.save(img_buffer, format="PNG")
                img_buffer.seek(0)

                overview_slide = prs.slides.add_slide(blank)
                header = overview_slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
                header.text_frame.text = "Dashboard Overview"
                header.text_frame.paragraphs[0].font.size = Pt(24)
                header.text_frame.paragraphs[0].font.bold = True

                aspect = image.height / image.width
                pic_width = Inches(12.3)
                pic_height = Inches(12.3 * aspect)
                max_height = Inches(6.6)
                if pic_height > max_height:
                    pic_height = max_height
                    pic_width = Inches(6.6 / aspect)

                overview_slide.shapes.add_picture(
                    img_buffer, Inches(0.5), Inches(0.85), width=pic_width, height=pic_height
                )
            except Exception as ex:
                logger.warning("Could not add dashboard overview slide: %s", ex)

            for widget in dashboard.widgets:
                if widget.chart is None:
                    continue
                try:
                    png_bytes = widget.chart.to_image(format="png", width=1200, height=700, scale=2)
                    slide = prs.slides.add_slide(blank)

                    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
                    header.text_frame.text = widget.title
                    header.text_frame.paragraphs[0].font.size = Pt(24)
                    header.text_frame.paragraphs[0].font.bold = True

                    image_stream = io.BytesIO(png_bytes)
                    slide.shapes.add_picture(image_stream, Inches(0.8), Inches(1.0), width=Inches(11.7))
                except Exception as ex:
                    logger.warning("Skipping widget '%s' in PPT export: %s", widget.title, ex)

            buffer = io.BytesIO()
            prs.save(buffer)
            return buffer.getvalue()

        except Exception as ex:
            logger.exception("PowerPoint export failed: %s", ex)
            return None

    # ------------------------------------------------------
    # HTML - self-contained, themed, single-file dashboard
    # ------------------------------------------------------

    def export_html(self, dashboard: Dashboard) -> str:
        theme = dashboard.theme

        kpi_cards = "".join(
            f"""
            <div class="kpi-card">
                <div class="kpi-label">{k.get('label', '')}</div>
                <div class="kpi-value">{k.get('value', '')}</div>
            </div>
            """
            for k in dashboard.kpis
        )

        chart_divs = ""
        for widget in dashboard.widgets:
            if widget.chart is None:
                continue
            try:
                chart_html = widget.chart.to_html(
                    full_html=False, include_plotlyjs=False, config={"displaylogo": False}
                )
            except Exception:
                continue
            chart_divs += f"""
            <div class="chart-card">
                <h3>{widget.title}</h3>
                {chart_html}
            </div>
            """

        recs_html = "".join(f"<li>{r}</li>" for r in dashboard.recommendations)

        html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{dashboard.title}</title>
<script src="https://cdn.plot.ly/plotly-2.32.0.min.js"></script>
<style>
    body {{
        margin: 0;
        font-family: {theme.font_family};
        background: {theme.background_gradient};
        color: {theme.text_color};
    }}
    .header {{
        background: {theme.header_gradient};
        color: {theme.header_text_color};
        padding: 32px 40px;
    }}
    .header h1 {{ margin: 0; font-size: 32px; }}
    .header p {{ margin: 6px 0 0; opacity: 0.9; }}
    .kpi-row {{
        display: flex; flex-wrap: wrap; gap: 16px; padding: 24px 40px 0;
    }}
    .kpi-card {{
        background: {theme.card_background};
        border-radius: {theme.border_radius}px;
        padding: 18px 22px; min-width: 160px; flex: 1;
        box-shadow: 0 4px 14px rgba(0,0,0,0.08);
    }}
    .kpi-label {{ color: {theme.muted_text}; font-size: 13px; text-transform: uppercase; }}
    .kpi-value {{ font-size: 26px; font-weight: 700; color: {theme.primary}; }}
    .chart-grid {{
        display: grid; grid-template-columns: repeat(auto-fit, minmax(420px, 1fr));
        gap: 20px; padding: 24px 40px 40px;
    }}
    .chart-card {{
        background: {theme.card_background};
        border-radius: {theme.border_radius}px;
        padding: 16px; box-shadow: 0 4px 14px rgba(0,0,0,0.08);
    }}
    .chart-card h3 {{ margin: 4px 0 8px; font-size: 15px; color: {theme.text_color}; }}
    .summary {{ padding: 0 40px; }}
    .summary-card {{
        background: {theme.card_background}; border-radius: {theme.border_radius}px;
        padding: 20px; box-shadow: 0 4px 14px rgba(0,0,0,0.08);
    }}
</style>
</head>
<body>
    <div class="header">
        <h1>{dashboard.title}</h1>
        <p>{dashboard.subtitle}</p>
    </div>
    <div class="kpi-row">{kpi_cards}</div>
    <div class="summary">
        <div class="summary-card">
            <p>{dashboard.summary}</p>
            <ul>{recs_html}</ul>
        </div>
    </div>
    <div class="chart-grid">{chart_divs}</div>
</body>
</html>"""
        return html

    # ------------------------------------------------------
    # Helpers
    # ------------------------------------------------------

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> tuple[float, float, float]:
        try:
            hex_color = hex_color.lstrip("#")
            if len(hex_color) != 6:
                return (0.15, 0.35, 0.85)
            r = int(hex_color[0:2], 16) / 255
            g = int(hex_color[2:4], 16) / 255
            b = int(hex_color[4:6], 16) / 255
            return (r, g, b)
        except Exception:
            return (0.15, 0.35, 0.85)

    @staticmethod
    def _hex_to_rgb255(hex_color: str) -> tuple[int, int, int]:
        try:
            hex_color = (hex_color or "").lstrip("#")
            if len(hex_color) != 6:
                return (37, 99, 235)
            return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
        except Exception:
            return (37, 99, 235)

    @staticmethod
    def _rounded_card(draw, box, theme, radius: int = 14):
        from PIL import ImageColor

        fill = DashboardExporter._hex_to_rgb255(theme.card_background) \
            if not str(theme.card_background).startswith("rgba") else (255, 255, 255)
        outline = DashboardExporter._hex_to_rgb255(theme.border_color)
        draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=1)

    @staticmethod
    def _load_fonts():
        from PIL import ImageFont
        import matplotlib

        base = matplotlib.get_data_path() + "/fonts/ttf/"
        try:
            regular = ImageFont.truetype(base + "DejaVuSans.ttf", 16)
            bold = ImageFont.truetype(base + "DejaVuSans-Bold.ttf", 20)
            title = ImageFont.truetype(base + "DejaVuSans-Bold.ttf", 30)
            small = ImageFont.truetype(base + "DejaVuSans.ttf", 13)
        except Exception:
            regular = bold = title = small = ImageFont.load_default()
        return regular, bold, title, small

